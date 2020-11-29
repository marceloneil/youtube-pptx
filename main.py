#!/usr/bin/env python3

import ffmpeg
import os
import shutil
import sys
import youtube_dl

from os.path import join, isfile
from pptx import Presentation

def my_hook(d):
    if d['status'] == 'finished':
        title = d['filename'].split('.')[0]
        mypath = join(os.getcwd(), title)
        mypptx = '{}.pptx'.format(title)
        if (isfile(mypptx)):
            return
        os.makedirs(mypath)
        stream = ffmpeg.input(d['filename'])
        stream = ffmpeg.filter(stream, 'fps', fps=1)
        stream = ffmpeg.output(stream, '{}/thumb%d.png'.format(title))
        ffmpeg.run(stream)

        prs = Presentation()
        blank_slide = prs.slide_layouts[6]
        height = prs.slide_height
        width = prs.slide_width

        thumbs = os.listdir(mypath)
        i = 1
        myfile = join(mypath, 'thumb{}.png'.format(i))
        print(myfile)
        while (isfile(myfile)):
            slide = prs.slides.add_slide(blank_slide)
            pic = slide.shapes.add_picture(myfile, 0, 0, height=height, width=width)
            i += 1
            myfile = join(mypath, 'thumb{}.png'.format(i))

        prs.save(mypptx)

vid = sys.argv[1]

ydl_opts = {
    'outtmpl': '%(id)s.%(ext)s',
    'progress_hooks': [my_hook]
}
ydl = youtube_dl.YoutubeDL(ydl_opts)
ydl.download([vid])
